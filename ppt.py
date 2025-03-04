import os
import json
from pptx import Presentation
from pptx.util import Pt, Inches  # Inches for positioning images

# Paths
NEWS_JSON = os.path.join("news_files", "final_news.json")
TEMPLATE_PPT = "template.pptx"
OUTPUT_PPT = "final_presentation.pptx"

def load_articles(json_path):
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)

def split_summary_into_sentences(summary):
    """
    Splits the summary text at full stops (".") into sentences.
    Ensures each sentence ends with a period.
    """
    sentences = [s.strip() for s in summary.split(".") if s.strip()]
    return [s if s.endswith(".") else s + "." for s in sentences]

def main():
    articles = load_articles(NEWS_JSON)
    prs = Presentation(TEMPLATE_PPT)

    for i, article in enumerate(articles):
        if i >= len(prs.slides):
            print("More articles than template slides. Stopping.")
            break

        slide = prs.slides[i]
        title_placeholder = slide.shapes.title

        # Find a suitable text placeholder for the summary
        summary_placeholder = None
        if len(slide.placeholders) > 1:
            summary_placeholder = slide.placeholders[1]
        else:
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title_placeholder:
                    summary_placeholder = shape
                    break

        # Insert the article title
        if title_placeholder:
            title_placeholder.text = article["title"]
            title_para = title_placeholder.text_frame.paragraphs[0]
            title_para.font.size = Pt(28)
            title_para.font.bold = True

        # Insert the article summary as bullet points
        if summary_placeholder and summary_placeholder.has_text_frame:
            text_frame = summary_placeholder.text_frame
            text_frame.clear()
            sentences = split_summary_into_sentences(article["summary"])
            for sentence in sentences:
                p = text_frame.add_paragraph()
                p.text = sentence
                p.font.size = Pt(22)
                p.font.name = "Arial"
                p.level = 0  # Main bullet level

        # Image handling logic (Check for both JPG & PNG)
        image_jpg = os.path.join("news_files", f"image_{i}.jpg")
        image_png = os.path.join("news_files", f"image_{i}.png")

        image_path = None
        if os.path.exists(image_jpg):
            image_path = image_jpg
        elif os.path.exists(image_png):
            image_path = image_png
        else:
            print(f"⚠️ Image missing for slide {i+1}: Tried {image_jpg} and {image_png}")

        # Insert image if found
        if image_path:
            left, top, width, height = Inches(1.09), Inches(2.02), Inches(7), Inches(6.2)  # Adjusted positioning
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"✅ Inserted image for slide {i+1}: {image_path}")

        print(f"✅ Populated slide {i+1} with article: {article['title']}")

    prs.save(OUTPUT_PPT)
    print(f"\nSaved final presentation to {OUTPUT_PPT}")
    return OUTPUT_PPT

if __name__ == "__main__":
    main()
